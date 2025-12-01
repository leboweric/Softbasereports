"""
QBR PowerPoint Generator
Generates Quarterly Business Review presentations using python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
import logging
from datetime import datetime
import os

logger = logging.getLogger(__name__)


class PPTXGenerator:
    """Generates QBR PowerPoint presentations from template"""

    def __init__(self, template_path: str):
        """
        Initialize generator with template path

        Args:
            template_path: Path to the PowerPoint template file
        """
        if not os.path.exists(template_path):
            raise FileNotFoundError(f"Template not found: {template_path}")

        self.template_path = template_path
        self.prs = None

    def generate_qbr_presentation(self, data: dict, output_path: str) -> str:
        """
        Generate a QBR presentation from data

        Args:
            data: Dictionary containing all QBR data
            output_path: Path to save the generated presentation

        Returns:
            Path to the generated file
        """
        try:
            # Load the template
            self.prs = Presentation(self.template_path)

            # Process each slide
            self._process_slides(data)

            # Save the presentation
            self.prs.save(output_path)
            logger.info(f"QBR presentation saved to: {output_path}")

            return output_path

        except Exception as e:
            logger.error(f"Error generating QBR presentation: {str(e)}")
            raise

    def _process_slides(self, data: dict):
        """Process all slides and replace placeholders"""

        customer = data.get('customer', {})
        quarter = data.get('quarter', 'Q4 2025')
        date_range = data.get('date_range', {})
        fleet_overview = data.get('fleet_overview', {})
        fleet_health = data.get('fleet_health', {})
        service_performance = data.get('service_performance', {})
        service_costs = data.get('service_costs', {})
        parts_rentals = data.get('parts_rentals', {})
        value_delivered = data.get('value_delivered', {})
        recommendations = data.get('recommendations', [])
        business_priorities = data.get('business_priorities', [])
        action_items = data.get('action_items', [])

        # Slide placeholder mapping
        replacements = {
            # Customer info
            '[Customer Name]': customer.get('customer_name', 'Customer'),
            '[Customer]': customer.get('customer_name', 'Customer'),
            'Q[X] [Year]': quarter,
            '[Quarter]': quarter,

            # Date range
            '[Start Date]': date_range.get('start', ''),
            '[End Date]': date_range.get('end', ''),
            '[Date Range]': f"{date_range.get('start', '')} - {date_range.get('end', '')}",

            # Fleet Overview
            '[Total Units]': str(fleet_overview.get('total_units', 0)),
            '[Under Contract]': str(fleet_overview.get('under_contract', 0)),
            '[Avg Age]': f"{fleet_overview.get('avg_age', 0):.1f}",
            '[Service Calls]': str(fleet_overview.get('service_calls', 0)),

            # Fleet Health
            '[Good Units]': str(fleet_health.get('good', 0)),
            '[Monitor Units]': str(fleet_health.get('monitor', 0)),
            '[Replace Units]': str(fleet_health.get('replace', 0)),

            # Service Performance
            '[Total WOs]': str(service_performance.get('total_work_orders', 0)),
            '[PM Compliance]': self._format_percent(service_performance.get('pm_compliance', 0)),
            '[Avg Response]': f"{service_performance.get('avg_response_time', 0):.1f}",
            '[First Fix Rate]': self._format_percent(service_performance.get('first_time_fix_rate', 0)),

            # Costs
            '[Total Service Cost]': self._format_currency(service_costs.get('total_service_cost', 0)),
            '[Labor Cost]': self._format_currency(service_costs.get('labor_cost', 0)),
            '[Parts Cost]': self._format_currency(service_costs.get('parts_cost', 0)),
            '[Cost Per Unit]': self._format_currency(service_costs.get('cost_per_unit', 0)),

            # Parts & Rentals
            '[Parts Total]': self._format_currency(parts_rentals.get('parts_total', 0)),
            '[Rental Revenue]': self._format_currency(parts_rentals.get('rental_revenue', 0)),
            '[Rental Days]': str(parts_rentals.get('rental_days', 0)),

            # Value Delivered
            '[PM Savings]': self._format_currency(value_delivered.get('pm_savings', 0)),
            '[Uptime Value]': self._format_currency(value_delivered.get('uptime_value', 0)),
            '[Total Value]': self._format_currency(value_delivered.get('total_value', 0)),

            # Generic placeholders
            '[##]': '0',
            '##': '0',
        }

        # Process each slide
        for slide in self.prs.slides:
            self._replace_placeholders_in_slide(slide, replacements)

            # Handle specific slide content based on slide index
            slide_idx = self.prs.slides.index(slide)

            # Try to identify slide by title
            title = self._get_slide_title(slide)

            if title:
                title_lower = title.lower()

                if 'business priorities' in title_lower:
                    self._populate_business_priorities(slide, business_priorities)
                elif 'recommendation' in title_lower:
                    self._populate_recommendations(slide, recommendations)
                elif 'action' in title_lower:
                    self._populate_action_items(slide, action_items)

    def _replace_placeholders_in_slide(self, slide, replacements: dict):
        """Replace all placeholders in a slide"""

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text = run.text
                        for placeholder, value in replacements.items():
                            if placeholder in text:
                                text = text.replace(placeholder, str(value))
                        run.text = text

            # Handle tables
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text_frame:
                            for paragraph in cell.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    text = run.text
                                    for placeholder, value in replacements.items():
                                        if placeholder in text:
                                            text = text.replace(placeholder, str(value))
                                    run.text = text

    def _get_slide_title(self, slide) -> str:
        """Get the title of a slide"""
        for shape in slide.shapes:
            if shape.has_text_frame:
                if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
                    if shape.placeholder_format.type == 1:  # Title placeholder
                        return shape.text
                # Also check for title-like shapes
                if shape.text and len(shape.text) < 100:
                    text = shape.text.strip()
                    if text and not any(c in text for c in ['[', ']', '$']):
                        return text
        return ""

    def _populate_business_priorities(self, slide, priorities: list):
        """Populate business priorities on a slide"""
        if not priorities:
            return

        # Find content placeholder or text frame
        for shape in slide.shapes:
            if shape.has_text_frame and not self._is_title_shape(shape):
                # Clear and populate
                tf = shape.text_frame
                # Keep first paragraph, clear rest
                while len(tf.paragraphs) > 1:
                    p = tf.paragraphs[-1]
                    p.clear()

                for i, priority in enumerate(priorities[:3]):
                    if i == 0:
                        tf.paragraphs[0].text = f"1. {priority.get('title', '')}"
                        if priority.get('description'):
                            tf.paragraphs[0].text += f"\n   {priority.get('description', '')}"
                    else:
                        p = tf.add_paragraph()
                        p.text = f"{i+1}. {priority.get('title', '')}"
                        if priority.get('description'):
                            p.text += f"\n   {priority.get('description', '')}"
                break

    def _populate_recommendations(self, slide, recommendations: list):
        """Populate recommendations on a slide"""
        if not recommendations:
            return

        # Find content placeholder
        for shape in slide.shapes:
            if shape.has_text_frame and not self._is_title_shape(shape):
                tf = shape.text_frame

                # Clear existing content
                for p in tf.paragraphs:
                    p.clear()

                for i, rec in enumerate(recommendations[:5]):
                    if i == 0:
                        tf.paragraphs[0].text = f"• {rec.get('title', '')}"
                    else:
                        p = tf.add_paragraph()
                        p.text = f"• {rec.get('title', '')}"

                    if rec.get('estimated_impact'):
                        p = tf.add_paragraph()
                        p.text = f"  Impact: {rec.get('estimated_impact', '')}"
                        p.level = 1
                break

    def _populate_action_items(self, slide, action_items: list):
        """Populate action items on a slide"""
        if not action_items:
            return

        # Look for a table to populate
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                # Assuming table has columns: Party, Action, Owner, Due Date
                for i, item in enumerate(action_items[:6]):
                    if i + 1 < len(table.rows):
                        row = table.rows[i + 1]  # Skip header row
                        if len(row.cells) >= 4:
                            row.cells[0].text = item.get('party', '')
                            row.cells[1].text = item.get('description', '')
                            row.cells[2].text = item.get('owner_name', '')
                            row.cells[3].text = item.get('due_date', '')
                break

    def _is_title_shape(self, shape) -> bool:
        """Check if a shape is a title placeholder"""
        if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
            return shape.placeholder_format.type == 1
        return False

    def _format_currency(self, value) -> str:
        """Format a value as currency"""
        if value is None:
            return '$0'
        try:
            return f"${value:,.0f}"
        except (ValueError, TypeError):
            return '$0'

    def _format_percent(self, value) -> str:
        """Format a value as percentage"""
        if value is None:
            return '0%'
        try:
            return f"{value * 100:.1f}%"
        except (ValueError, TypeError):
            return '0%'


def generate_qbr_pptx(template_path: str, data: dict, output_path: str) -> str:
    """
    Convenience function to generate QBR PowerPoint

    Args:
        template_path: Path to template file
        data: QBR data dictionary
        output_path: Where to save the generated file

    Returns:
        Path to generated file
    """
    generator = PPTXGenerator(template_path)
    return generator.generate_qbr_presentation(data, output_path)
